"""
PowerPoint converter for DocBridge
"""
from typing import Optional, Union, List
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

from .base import BaseConverter
from ..styles.default import DefaultStyle


class PowerPointConverter(BaseConverter):
    """
    Convert Markdown to PowerPoint presentations.
    """

    def __init__(self, style: Optional[DefaultStyle] = None):
        super().__init__(style)
        self.prs = None

    def convert(self, markdown_text: str) -> Presentation:
        """
        Convert markdown text to PowerPoint presentation.

        Args:
            markdown_text: The markdown text to convert

        Returns:
            pptx Presentation object
        """
        # Create presentation
        prs = Presentation()

        # Set slide dimensions
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        # Parse markdown and create slides
        slides_content = self._parse_slides(markdown_text)

        for slide_content in slides_content:
            self._add_slide(prs, slide_content)

        self._document = prs
        return prs

    def convert_file(self, input_path: Union[str, Path],
                     output_path: Optional[Union[str, Path]] = None) -> str:
        """
        Convert markdown file to PowerPoint file.

        Args:
            input_path: Path to input markdown file
            output_path: Path for output file (optional)

        Returns:
            Path to output file
        """
        # Read input file
        markdown_text = self._read_file(input_path)

        # Convert to presentation
        prs = self.convert(markdown_text)

        # Determine output path
        output_path = self._ensure_output_path(input_path, output_path, '.pptx')

        # Save presentation
        prs.save(output_path)

        return str(output_path)

    def _parse_slides(self, text: str) -> List[dict]:
        """
        Parse markdown into slide content.

        Each H1 or H2 starts a new slide.
        """
        lines = text.split('\n')
        slides = []
        current_slide = {'title': '', 'content': []}

        for line in lines:
            # Check for heading (new slide)
            if line.startswith('# ') or line.startswith('## '):
                # Save previous slide if it has content
                if current_slide['title'] or current_slide['content']:
                    slides.append(current_slide)

                # Start new slide
                level = 1 if line.startswith('# ') else 2
                title = line.lstrip('#').strip()
                current_slide = {
                    'title': title,
                    'level': level,
                    'content': []
                }
            elif line.strip():
                current_slide['content'].append(line)

        # Don't forget last slide
        if current_slide['title'] or current_slide['content']:
            slides.append(current_slide)

        # If no slides were created, create one with all content
        if not slides:
            slides = [{
                'title': 'Presentation',
                'level': 1,
                'content': lines
            }]

        return slides

    def _add_slide(self, prs: Presentation, slide_content: dict) -> None:
        """Add a slide to presentation."""
        # Choose layout based on content
        if slide_content['level'] == 1:
            # Title slide layout
            slide_layout = prs.slide_layouts[0]  # Title slide
        else:
            # Content slide layout
            slide_layout = prs.slide_layouts[1]  # Title and content

        slide = prs.slides.add_slide(slide_layout)

        # Set title
        if slide.shapes.title:
            title = slide.shapes.title
            title.text = slide_content['title']

            # Style title
            for paragraph in title.text_frame.paragraphs:
                paragraph.font.size = Pt(32 if slide_content['level'] == 1 else 28)
                paragraph.font.bold = True
                paragraph.alignment = PP_ALIGN.LEFT

        # Add content
        if slide_content['content'] and len(slide.placeholders) > 1:
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame
            tf.clear()

            for i, line in enumerate(slide_content['content']):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()

                # Process line
                text = line.strip()

                # Remove markdown formatting for now
                text = text.lstrip('- ').lstrip('* ').lstrip('1. ')

                p.text = text
                p.font.size = Pt(18)
                p.level = 0

                # Check for list items
                if line.strip().startswith(('- ', '* ')):
                    p.level = 0
                elif line.strip().startswith(('  - ', '  * ')):
                    p.level = 1
